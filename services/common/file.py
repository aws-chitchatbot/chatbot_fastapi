import csv
import os
from io import BufferedReader
from typing import Optional
from PyPDF2 import PdfReader
from fastapi import UploadFile
import mimetypes
from loguru import logger
import docx2txt
import csv
import pptx
from models.models import Document, DocumentMetadata
import re
from datetime import datetime


async def get_document_from_file(
    file: UploadFile, metadata: DocumentMetadata
) -> Document:
    """추출한 텍스트를 기반으로 Document를 반환한다"""
    extracted_text = await extract_text_from_form_file(file)
    document = Document(text=extracted_text, metadata=metadata)

    return document


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """파일 경로가 지정된 파일의 텍스트 내용을 반환"""
    print(mimetype)
    if mimetype is None:
        # 확장자를 기반으로 파일의 mimetype을 가져온다
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("지원하지 않는 파일 형식입니다.")

    try:
        with open(filepath, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(e)
        raise e
    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # 텍스트 파일로 부터 텍스트를 읽기
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # 지원하지 않는 파일 형식
        raise ValueError("Unsupported file type: {}".format(mimetype))
    return extracted_text


# mimetype에 따라 파일에서 텍스트 추출
async def extract_text_from_form_file(file: UploadFile):
    """파일의 텍스트 내용을 반환"""
    # 업로드 파일 개체에서 파일 내용 가져오기
    mimetype = file.content_type
    logger.info(f"mimetype: {mimetype}")
    logger.info(f"file.file: {file.file}")
    logger.info("file: ", file)

    file_stream = await file.read()
    temp_file_path = "/tmp/temp_file"

    # 파일을 임시 위치에 쓰기
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)
    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        logger.error(e)
        os.remove(temp_file_path)
        raise e

    # 임시 파일 제거
    os.remove(temp_file_path)

    return extracted_text


def extract_messages(input_file, target_people):
    # 파일 읽기
    with open(input_file, "r", encoding="utf-8") as file:
        content = file.read()

    # 날짜 추출
    date_match = re.search(r"\[.*\] (\d{4}-\d{2}-\d{2})", content)
    if date_match:
        date_str = date_match.group(1)
    else:
        print("파일에서 날짜를 찾을수 없습니다.")
        return

    formatted_date = datetime.strptime(date_str, "%Y-%m-%d").strftime("%Y-%m-%d")

    # 정규표현식
    # 여러 사람에 대한 이름을 OR(|) 연산자를 사용하여 정규표현식으로 만든다
    target_person_pattern = "|".join(map(re.escape, target_people))
    pattern = re.compile(
        rf"\[({target_person_pattern})\] \[.*?\] (.*?)(?=\n\[.*?\]|\Z)", re.DOTALL
    )
    matches = re.findall(pattern, content)

    # 추출
    output_folder = "text"
    os.makedirs(output_folder, exist_ok=True)
    output_file = os.path.join(output_folder, f"{formatted_date}.txt")

    with open(output_file, "w", encoding="utf-8") as output:
        for match in matches:
            output.write(f"{match[0]}: {match[1].strip()}\n")

    print(f"다음 이름으로 저장되었습니다.: {output_file}")
